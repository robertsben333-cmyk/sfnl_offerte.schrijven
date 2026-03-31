"""PPTX aanpak_overview slide — chevron phases with descriptions."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    find_shape,
    hex_color as _hex,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add aanpak overview slide with 2-5 phase chevrons.

    content keys:
      title (str): Slide title e.g. "ONZE AANPAK"
      subtitle (str): One-line summary
      phases (list): Each has naam, beschrijving, tijdlijn
      proposition (str): Proposition id for accent color
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "aanpak_overview")
    primary = _hex("primary")
    muted = _hex("text_muted")
    style = content.get("__style_context__", {}).get("aanpak_overview", {})

    set_text_preserve(find_placeholder(slide, 0), content.get("title", "PLAN VAN AANPAK"))
    set_text_preserve(find_placeholder(slide, 1), content.get("subtitle", ""))

    phases = content.get("phases", [])[:4]
    if not phases:
        return

    phase_shapes = [
        find_shape(slide, "Arrow: Pentagon 5"),
        find_shape(slide, "Arrow: Chevron 6"),
        find_shape(slide, "Arrow: Chevron 7"),
        find_shape(slide, "Arrow: Chevron 8"),
    ]
    desc_shapes = [
        find_shape(slide, "Tijdelijke aanduiding voor inhoud 19"),
        find_shape(slide, "Tijdelijke aanduiding voor inhoud 20"),
        find_shape(slide, "Tijdelijke aanduiding voor inhoud 23"),
        find_shape(slide, "Tijdelijke aanduiding voor inhoud 26"),
    ]
    intro_shape = find_shape(slide, "Tijdelijke aanduiding voor inhoud 18")
    set_paragraphs(
        intro_shape,
        [{
            "text": content.get("intro", ""),
            "role": "body",
            "size": style.get("intro_size", 12),
            "color": primary,
        }],
    )

    timeline_shapes = [
        find_shape(slide, "TextBox 14"),
        find_shape(slide, "TextBox 16"),
        find_shape(slide, "TextBox 17", 0),
    ]
    year_shape = find_shape(slide, "TextBox 17", 1)
    disclaimer_shape = find_shape(slide, "Tekstvak 1")

    for idx in range(4):
        phase = phases[idx] if idx < len(phases) else {}
        set_text_preserve(phase_shapes[idx], phase.get("naam", ""))
        set_paragraphs(
            desc_shapes[idx],
            [{
                "text": phase.get("beschrijving", ""),
                "role": "body",
                "size": style.get("desc_size", 11),
                "color": primary,
            }],
        )

    for idx in range(3):
        label = phases[idx].get("tijdlijn", "") if idx < len(phases) else ""
        set_paragraphs(
            timeline_shapes[idx],
            [{
                "text": label,
                "role": "body",
                "size": style.get("timeline_size", 9),
                "color": muted,
            }],
        )

    year_text = content.get("year", "")
    if not year_text and phases:
        combined = " ".join(phase.get("tijdlijn", "") for phase in phases)
        for token in combined.split():
            if token.isdigit() and len(token) == 4:
                year_text = token
                break
    set_paragraphs(
        year_shape,
        [{
            "text": year_text,
            "role": "body",
            "size": style.get("timeline_size", 9),
            "color": muted,
        }],
    )
    set_paragraphs(
        disclaimer_shape,
        [{
            "text": content.get("disclaimer", "* Planning is indicatief. De planning wordt definitief gemaakt na ondertekening van de offerte"),
            "role": "body",
            "size": style.get("timeline_size", 9),
            "italic": True,
            "color": muted,
        }],
    )
