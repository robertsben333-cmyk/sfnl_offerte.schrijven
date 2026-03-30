"""PPTX two-column slide for free-form proposition content."""
from pptx import Presentation
from pptx.util import Pt
from skills.pptx_offerte.scripts.slides._utils import STYLE, ACCENT_MAP, hex_color as _hex, blank_layout


def _add_text_block(slide, left, top, width, height, heading, body, accent, primary):
    if heading:
        heading_box = slide.shapes.add_textbox(left, top, width, int(height * 0.12))
        paragraph = heading_box.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = heading
        run.font.name = STYLE["fonts"]["heading"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = accent
        top += int(height * 0.12)
        height -= int(height * 0.12)

    if body:
        body_box = slide.shapes.add_textbox(left, top, width, height)
        body_box.text_frame.word_wrap = True
        for index, item in enumerate(body if isinstance(body, list) else [body]):
            paragraph = (
                body_box.text_frame.paragraphs[0]
                if index == 0
                else body_box.text_frame.add_paragraph()
            )
            run = paragraph.add_run()
            run.text = item
            run.font.name = STYLE["fonts"]["body"]
            run.font.size = Pt(10)
            run.font.color.rgb = primary


def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add a free-form two-column slide.

    content keys:
      title, subtitle, left_title, left_body, right_title, right_body, proposition
    """
    W = prs.slide_width
    H = prs.slide_height
    slide = prs.slides.add_slide(blank_layout(prs))

    proposition = content.get("proposition", "mbc")
    accent = _hex(ACCENT_MAP.get(proposition, "accent_mbc"))
    primary = _hex("primary")
    white = _hex("white")

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = white

    title = content.get("title", "")
    if title:
        tf = slide.shapes.add_textbox(int(W * 0.036), int(H * 0.08), int(W * 0.939), int(H * 0.049))
        paragraph = tf.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = title
        run.font.name = STYLE["fonts"]["heading"]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = accent

    subtitle = content.get("subtitle", "")
    if subtitle:
        tf = slide.shapes.add_textbox(int(W * 0.036), int(H * 0.138), int(W * 0.939), int(H * 0.083))
        tf.text_frame.word_wrap = True
        paragraph = tf.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = subtitle
        run.font.name = STYLE["fonts"]["subtitle"]
        run.font.size = Pt(11)
        run.font.color.rgb = primary

    col_top = int(H * 0.26)
    col_h = int(H * 0.62)
    gap = int(W * 0.03)
    col_w = int((W * 0.939 - gap) / 2)
    left = int(W * 0.036)
    right = left + col_w + gap

    _add_text_block(
        slide,
        left,
        col_top,
        col_w,
        col_h,
        content.get("left_title", ""),
        content.get("left_body", ""),
        accent,
        primary,
    )
    _add_text_block(
        slide,
        right,
        col_top,
        col_w,
        col_h,
        content.get("right_title", ""),
        content.get("right_body", ""),
        accent,
        primary,
    )
