"""PPTX team slide — 2 or 3 team members."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

STYLE_PATH = os.path.join(os.path.dirname(__file__), "../../../../data/style.json")
with open(STYLE_PATH) as f:
    STYLE = json.load(f)

ACCENT_MAP = {
    "mbc": "accent_mbc",
    "impact_meten": "accent_impact_meten",
    "advies_innovatieve_financiering": "accent_advies_innovatief",
    "intermediair": "accent_intermediair",
    "fondsmanagement": "accent_fondsmanagement",
    "partnerschappen": "accent_partnerschappen",
}

def _hex(key: str) -> RGBColor:
    h = STYLE["colors"][key].lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add team slide with 2-3 members.
    content: members (list of {name, title, bio}), proposition
    """
    W = prs.slide_width
    H = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    proposition = content.get("proposition", "mbc")
    accent = _hex(ACCENT_MAP.get(proposition, "accent_mbc"))
    primary = _hex("primary")
    white = _hex("white")

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = white

    # Slide title
    tf = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), W - Inches(1.0), Inches(0.5))
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "ONS TEAM"
    run.font.name = STYLE["fonts"]["heading"]
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = accent

    members = content.get("members", [])[:3]
    if not members:
        return

    n = len(members)
    col_w = int((W - Inches(1.0)) / n)
    top = int(H * 0.22)
    card_h = int(H * 0.65)

    for i, member in enumerate(members):
        left = int(Inches(0.5) + i * (col_w + Inches(0.05)))

        # Accent bar
        bar = slide.shapes.add_shape(1, left, top, col_w - Inches(0.1), int(Inches(0.06)))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Name
        tf = slide.shapes.add_textbox(left, int(top + Inches(0.1)), col_w - Inches(0.1), Inches(0.45))
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = member.get("name", "")
        run.font.name = STYLE["fonts"]["heading"]
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = primary

        # Title
        tf = slide.shapes.add_textbox(left, int(top + Inches(0.58)), col_w - Inches(0.1), Inches(0.35))
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = member.get("title", "")
        run.font.name = STYLE["fonts"]["body"]
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = accent

        # Bio
        tf = slide.shapes.add_textbox(left, int(top + Inches(0.98)), col_w - Inches(0.1), int(card_h - Inches(0.98)))
        tf.text_frame.word_wrap = True
        p = tf.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = member.get("bio", "")
        run.font.name = STYLE["fonts"]["body"]
        run.font.size = Pt(10)
        run.font.color.rgb = primary
