"""PPTX tijdslijn slide — preserve the template schedule table and fills."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    set_lines_preserve,
    set_text_preserve,
)

_SUMMARY_ROWS = [1, 5, 10, 13]
_MONTH_COLS = [1, 2, 3, 4, 5]
_DETAIL_ROWS = [
    [2, 3, 4],
    [6, 7, 8, 9],
    [11, 12],
    [14],
]


def _header_labels(content: dict, phases: list[dict]) -> list[str]:
    labels = content.get("timeline_labels")
    if labels:
        return list(labels)[:5]
    return [phase.get("periode", "") for phase in phases[:5]]


def _display_phases(content: dict) -> list[dict]:
    if content.get("display_phases"):
        return list(content["display_phases"])[:4]

    phases = []
    for idx, phase in enumerate(content.get("phases", [])[:4], start=1):
        phases.append({
            "naam": phase.get("naam", f"Fase {idx}"),
            "periode": phase.get("periode", ""),
            "dagen": phase.get("dagen", phase.get("inzet", "")),
            "slots": phase.get("slots", [min(idx, 5)]),
        })
    return phases


def _phase_steps(display_phase: dict, raw_phase: dict, max_rows: int) -> list[str]:
    steps = display_phase.get("substeps") or raw_phase.get("substeps") or []
    labels: list[str] = []
    for step in steps[:max_rows]:
        if isinstance(step, dict):
            label = str(step.get("label", "")).strip()
        else:
            label = str(step).strip()
        if label:
            labels.append(label)
    if not labels and raw_phase.get("activiteiten"):
        labels = [str(raw_phase["activiteiten"]).strip()]
    return labels[:max_rows]


def add_slide(prs: Presentation, content: dict) -> None:
    """content: title, intro, disclaimer, phases/display_phases, timeline_labels"""
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "tijdslijn")

    set_text_preserve(find_placeholder(slide, 0), content.get("title", "TIJDSLIJN"))

    subtitle_lines = []
    if content.get("intro"):
        subtitle_lines.append(content["intro"])
    if content.get("disclaimer"):
        subtitle_lines.append(f"* {content['disclaimer']}")
    if subtitle_lines:
        set_lines_preserve(find_placeholder(slide, 1), subtitle_lines)

    table = next(shape for shape in slide.shapes if shape.has_table).table

    labels = _header_labels(content, content.get("phases", []))
    start_year = content.get("start_year", "")
    end_year = content.get("end_year", "")
    if not start_year:
        for label in labels:
            for token in str(label).replace("–", "-").split():
                if token.isdigit() and len(token) == 4:
                    start_year = token
                    break
            if start_year:
                break
    set_text_preserve(table.cell(0, 0), start_year)
    for col_idx, label in zip(_MONTH_COLS, labels):
        set_text_preserve(table.cell(0, col_idx), str(label))
    for col_idx in _MONTH_COLS[len(labels):]:
        set_text_preserve(table.cell(0, col_idx), "")
    set_text_preserve(table.cell(0, 6), end_year)
    set_text_preserve(table.cell(0, 7), content.get("effort_label", "Inzet"))

    phases = _display_phases(content)
    for idx, summary_row in enumerate(_SUMMARY_ROWS):
        phase = phases[idx] if idx < len(phases) else {}
        set_text_preserve(table.cell(summary_row, 0), phase.get("naam", ""))
        days = phase.get("dagen", "")
        days_text = f"{days} dagen" if days not in ("", None) else ""
        set_text_preserve(table.cell(summary_row, 7), days_text)

    raw_phases = list(content.get("phases", []))
    for phase_idx, detail_rows in enumerate(_DETAIL_ROWS):
        display_phase = phases[phase_idx] if phase_idx < len(phases) else {}
        raw_phase = raw_phases[phase_idx] if phase_idx < len(raw_phases) else {}
        detail_labels = _phase_steps(display_phase, raw_phase, len(detail_rows))
        for row_idx, label in zip(detail_rows, detail_labels):
            set_text_preserve(table.cell(row_idx, 0), label)
        for row_idx in detail_rows[len(detail_labels):]:
            set_text_preserve(table.cell(row_idx, 0), "")
