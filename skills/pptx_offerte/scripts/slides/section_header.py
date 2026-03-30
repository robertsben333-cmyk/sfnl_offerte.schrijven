"""PPTX section header (divider) slide component."""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

STYLE_PATH = os.path.join(os.path.dirname(__file__), "../../../../data/style.json")
with open(STYLE_PATH) as f:
    STYLE = json.load(f)

ACCENT_MAP = {
    "mbc": "accent_mbc",
    "impact_meten": "accent_impact_meten",
    "advies_innovatieve_financiering": "accent_advies_innovatief",
    "intermediair": "accent_intermediair",
    "fondsmanagement": "accent_fondsmanagement",
    "partnerschappen": "accent_partnerschappen",
}

def _hex(key: str) -> RGBColor:
    h = STYLE["colors"][key].lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add a section divider slide with colored background.
    content: title, proposition
    """
    W = prs.slide_width
    H = prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    proposition = content.get("proposition", "mbc")
    accent = _hex(ACCENT_MAP.get(proposition, "accent_mbc"))
    white = _hex("white")

    # Full colored background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = accent

    # Section title
    tf = slide.shapes.add_textbox(Inches(1), int(H * 0.35), W - Inches(2), Inches(1.2))
    tf.text_frame.word_wrap = True
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = content.get("title", "")
    run.font.name = STYLE["fonts"]["heading"]
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = white
