"""PPTX randvoorwaarden slide — fill the existing template placeholders."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """content: title, subtitle, items (list[str])"""
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "randvoorwaarden")
    style = content.get("__style_context__", {}).get("randvoorwaarden", {})

    set_text_preserve(find_placeholder(slide, 0), content.get("title", "RANDVOORWAARDEN VOOR SUCCES"))
    if content.get("subtitle"):
        set_text_preserve(find_placeholder(slide, 1), content["subtitle"])

    # Use fewer, stronger bullets on this slide.
    lines = [item for item in content.get("items", [])[:4]]
    shape = find_placeholder(slide, 10)
    paragraphs = [{"text": f"• {item}", "role": "body", "size": style.get("body_size", 13)} for item in lines]
    set_paragraphs(shape, paragraphs)
