"""PPTX aanleiding slide — vraagstuk, uitdagingen, behoefte."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    ACCENT_MAP,
    clone_template_slide,
    find_placeholder,
    find_shape,
    hex_color as _hex,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add aanleiding slide with three text blocks.
    content: summary_line, vraagstuk, uitdagingen, behoefte, proposition
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "aanleiding")

    proposition = content.get("proposition", "mbc")
    accent = _hex(ACCENT_MAP.get(proposition, "accent_mbc"))
    muted = _hex("text_muted")
    style = content.get("__style_context__", {}).get("aanleiding", {})

    set_paragraphs(
        find_placeholder(slide, 0),
        [{
            "text": content.get("title", "AANLEIDING OFFERTE"),
            "role": "heading",
            "size": 20,
            "bold": True,
            "color": accent,
        }],
    )
    if content.get("summary_line"):
        set_paragraphs(
            find_placeholder(slide, 1),
            [{
                "text": content["summary_line"],
                "role": "subtitle",
                "size": style.get("summary_size", 13),
                "color": muted,
            }],
        )

    white = _hex("white")

    def _fill_block(shape_name: str, body: str) -> None:
        panel = find_shape(slide, shape_name)
        body_box = slide.shapes.add_textbox(
            panel.left + int(panel.width * 0.02),
            panel.top + int(panel.height * 0.16),
            int(panel.width * 0.96),
            int(panel.height * 0.74),
        )
        body_box.text_frame.word_wrap = True
        body_box.fill.background()
        body_box.line.fill.background()
        set_paragraphs(
            body_box,
            [{"text": body, "role": "body", "size": style.get("body_size", 12), "color": white}],
        )

    set_text_preserve(find_shape(slide, "Rectangle 1"), "Maatschappelijk vraagstuk")
    set_text_preserve(find_shape(slide, "Rectangle 8"), "Grootste uitdagingen")
    set_text_preserve(find_shape(slide, "Rectangle 9"), "Behoefte van de klant")
    _fill_block("Rectangle 1", content.get("vraagstuk", ""))
    _fill_block("Rectangle 8", content.get("uitdagingen", ""))
    _fill_block("Rectangle 9", content.get("behoefte", ""))
