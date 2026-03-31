"""PPTX budget_table slide — use the existing template table and notes styling."""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    hex_color as _hex,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    content: rows (list of {fase, dagen, kosten}), day_rate, tarief_note,
             social_rate_disclaimer, termijnen, proposition
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "budget_table")
    style = content.get("__style_context__", {}).get("budget_table", {})
    white = _hex("white")
    primary = _hex("primary")

    set_text_preserve(find_placeholder(slide, 0), content.get("title", "BEGROTING"))
    if content.get("subtitle"):
        set_text_preserve(find_placeholder(slide, 1), content["subtitle"])

    rows = content.get("rows", [])
    table = next(shape for shape in slide.shapes if shape.has_table).table

    def _set_cell(cell, text: str, *, bold: bool = False, color=None, alignment=PP_ALIGN.LEFT) -> None:
        set_paragraphs(cell, [{
            "text": text,
            "role": "body",
            "size": 10,
            "bold": bold,
            "color": color,
            "alignment": alignment,
        }])

    headers = ["", "# dagen", "Dagtarief ex. btw", "Kosten"]
    for col_idx, header in enumerate(headers):
        _set_cell(table.cell(0, col_idx), header, bold=True, color=white, alignment=PP_ALIGN.CENTER)

    max_body_rows = len(table.rows) - 1
    for row_idx in range(1, len(table.rows)):
        for col_idx in range(len(table.columns)):
            _set_cell(table.cell(row_idx, col_idx), "")

    default_day_rate = content.get("day_rate", "")
    display_rows = rows[:max_body_rows]
    if len(rows) > max_body_rows and max_body_rows > 0:
        display_rows = rows[: max_body_rows - 1] + [rows[-1]]

    for row_idx, row in enumerate(display_rows, start=1):
        day_rate = row.get("dagtarief", default_day_rate)
        day_rate_text = f"\u20ac {day_rate:,.0f}".replace(",", ".") if isinstance(day_rate, (int, float)) else str(day_rate)
        kosten = row.get("kosten", "")
        kosten_text = f"\u20ac {kosten:,.0f}".replace(",", ".") if isinstance(kosten, (int, float)) else str(kosten)
        values = [row.get("fase", ""), str(row.get("dagen", "")), day_rate_text, kosten_text]
        is_total = "totaal" in str(row.get("fase", "")).lower()
        for col_idx, value in enumerate(values):
            align = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
            _set_cell(table.cell(row_idx, col_idx), value, bold=is_total, color=primary, alignment=align)

    tarief_note = content.get("tarief_note", "")
    social_rate_disclaimer = content.get("social_rate_disclaimer", "")
    termijnen = content.get("termijnen", [])
    notes = [
        "Social Finance NL voert de opdracht uit op basis van onze Algemene Voorwaarden.",
        "Het tarief is exclusief BTW en op basis van 8 uur per dag.",
        "Het tarief is een teamtarief gebaseerd op een team bestaande uit een director, manager en associate/analyst.",
    ]
    if tarief_note:
        notes.append(tarief_note)
    if social_rate_disclaimer:
        notes.append(social_rate_disclaimer)
    if termijnen:
        notes.append(f"Facturatie geschiedt volgens het volgende schema: {'; '.join(termijnen)}.")
    notes.append("Zonder vooraf verkregen toestemming zal SFNL geen ‘meerwerk’ in rekening brengen. Als we over het aantal uren heen gaan is dat het risico van Social Finance NL tenzij de opdracht is uitgebreid.")
    shape = find_placeholder(slide, 10)
    set_paragraphs(shape, [{"text": f"• {note}", "role": "body", "size": style.get("notes_size", 10)} for note in notes])
