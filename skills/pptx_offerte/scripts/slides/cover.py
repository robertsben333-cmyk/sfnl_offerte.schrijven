"""PPTX cover slide component."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    find_shape,
    hex_color as _hex,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add a cover slide.
    content: client, title, date, proposition
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "cover")
    white = _hex("white")

    set_text_preserve(find_shape(slide, "Text Placeholder 4"), content.get("title", ""))

    meta = []
    if content.get("client"):
        meta.append({
            "text": content["client"],
            "role": "subtitle",
            "size": 20,
            "bold": True,
            "color": white,
        })
    if content.get("date"):
        meta.append({
            "text": content["date"],
            "role": "body",
            "size": 18,
            "bold": True,
            "color": white,
        })
    if meta:
        set_paragraphs(find_placeholder(slide, 14), meta)
