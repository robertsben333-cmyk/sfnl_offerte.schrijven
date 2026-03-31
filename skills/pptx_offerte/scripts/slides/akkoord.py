"""PPTX akkoord slide — fill the existing template text boxes and placeholders."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clear_run_highlights,
    clone_template_slide,
    find_placeholder,
    find_shape,
    set_paragraphs,
    set_lines_preserve,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    content: title, subtitle, randvoorwaarden_tekst, termijnen (list[str]),
             sfnl_naam, klant_naam, klant_org
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "akkoord")
    style = content.get("__style_context__", {}).get("akkoord", {})

    set_text_preserve(find_placeholder(slide, 0), content.get("title", "RANDVOORWAARDEN EN AKKOORD"))
    if content.get("subtitle"):
        set_text_preserve(find_placeholder(slide, 1), content["subtitle"])

    body_lines = [
        {"text": "Randvoorwaarden", "role": "body", "size": style.get("body_size", 12)},
        {"text": content.get("randvoorwaarden_tekst", ""), "role": "body", "size": style.get("body_size", 12)},
    ]
    body_lines.extend(
        {"text": termijn, "role": "body", "size": style.get("body_size", 12)}
        for termijn in content.get("termijnen", [])[:3]
    )
    set_paragraphs(find_placeholder(slide, 10), body_lines)

    set_lines_preserve(
        find_shape(slide, "TextBox 9"),
        ["", content.get("sfnl_naam", ""), "Social Finance NL"],
    )
    klant_box = find_shape(slide, "TextBox 8", 0)
    set_lines_preserve(
        klant_box,
        ["", content.get("klant_naam", ""), content.get("klant_org", "")],
    )
    clear_run_highlights(klant_box)
