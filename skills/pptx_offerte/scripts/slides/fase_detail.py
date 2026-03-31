"""PPTX fase_detail slide — two-column layout with actions left, doel/aanpak right."""
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import (
    clone_template_slide,
    find_placeholder,
    hex_color as _hex,
    set_paragraphs,
    set_text_preserve,
)


def add_slide(prs: Presentation, content: dict) -> None:
    """
    Add a fase detail slide.

    content keys:
      number (int): Phase number
      naam (str): Phase name
      klant (str): Client name for "Acties [klant]" label
      doel (str): Phase goal
      aanpak (str): Approach description
      acties_sfnl (list[str]): SFNL actions
      acties_klant (list[str]): Client actions
      deliverable (str): Deliverable description
      dagen (int): Number of days
      tijdlijn (str): Timeline string
      proposition (str): Proposition id
    """
    target_idx = content.get("__target_slide_index__")
    slide = prs.slides[target_idx] if target_idx is not None else clone_template_slide(prs, "fase_detail")
    primary = _hex("primary")
    style = content.get("__style_context__", {}).get("fase_detail", {})

    number = content.get("number", "")
    naam = content.get("naam", "")
    set_text_preserve(find_placeholder(slide, 0), f"{number}. {naam}" if number != "" else naam)
    set_text_preserve(find_placeholder(slide, 1), content.get("subtitle", f"AANPAK FASE {number}" if number != "" else ""))

    right_lines = [
        {"text": "Doel", "role": "body", "size": style.get("right_label_size", 13), "bold": True, "color": primary},
        {"text": content.get("doel", ""), "role": "body", "size": style.get("right_body_size", 13), "color": primary},
        {"text": "", "role": "body", "size": style.get("right_body_size", 13), "color": primary},
        {"text": "Aanpak", "role": "body", "size": style.get("right_label_size", 13), "bold": True, "color": primary},
        {"text": content.get("aanpak", ""), "role": "body", "size": style.get("right_body_size", 13), "color": primary},
        {"text": "", "role": "body", "size": style.get("right_body_size", 13), "color": primary},
        {"text": "", "role": "body", "size": style.get("right_body_size", 13), "color": primary},
        {"text": "", "role": "body", "size": style.get("right_body_size", 13), "color": primary},
    ]
    set_paragraphs(find_placeholder(slide, 11), right_lines)

    klant = content.get("klant", "Klant").upper()
    acties_sfnl = [item for item in content.get("acties_sfnl", [])[:3]]
    acties_klant = [item for item in content.get("acties_klant", [])[:3]]
    left_paragraphs = []

    def _append_section(label: str, lines: list[str]) -> None:
        nonlocal left_paragraphs
        clean_lines = [line for line in lines if str(line).strip()]
        if not clean_lines:
            return
        if left_paragraphs:
            left_paragraphs.append({
                "text": "",
                "role": "body",
                "size": style.get("left_body_size", 10),
                "color": primary,
            })
        left_paragraphs.append({
            "text": label,
            "role": "body",
            "size": style.get("left_label_size", 11),
            "bold": True,
            "color": primary,
        })
        for line in clean_lines:
            left_paragraphs.append({
                "text": f"• {line}" if not label.startswith("Duur") else line,
                "role": "body",
                "size": style.get("left_body_size", 10),
                "color": primary,
            })

    _append_section("Acties Social Finance NL", acties_sfnl)
    _append_section(f"Acties {klant}", acties_klant)
    _append_section("Deliverable", [content.get("deliverable", "")])
    _append_section(
        "Duur",
        [
            f"{content.get('dagen', '')} dagdelen" if content.get("dagen", "") != "" else "",
            str(content.get("tijdlijn", "")),
        ],
    )

    set_paragraphs(find_placeholder(slide, 12), left_paragraphs)
