"""Shared helpers for PPTX slide components."""
import json, os
from copy import deepcopy
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.oxml.ns import qn

STYLE_PATH = os.path.join(os.path.dirname(__file__), "../../../../data/style.json")
with open(STYLE_PATH) as f:
    STYLE = json.load(f)

ACCENT_MAP = {
    "mbc": "accent_mbc",
    "impact_meten": "accent_impact_meten",
    "advies_innovatieve_financiering": "accent_advies_innovatief",
    "intermediair": "accent_intermediair",
    "fondsmanagement": "accent_fondsmanagement",
    "partnerschappen": "accent_partnerschappen",
}

TEMPLATE_SLIDES = {
    "cover": 1,
    "aanleiding": 4,
    "aanpak_overview": 6,
    "fase_detail": 7,
    "tijdslijn": 8,
    "randvoorwaarden": 9,
    "budget_table": 13,
    "akkoord": 14,
}


def hex_color(key: str) -> RGBColor:
    h = STYLE["colors"][key].lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def blank_layout(prs: Presentation):
    """Return the blank slide layout by name, falling back to index 0."""
    name = STYLE.get("blank_layout_name", "Blank")
    for layout in prs.slide_master.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_master.slide_layouts[0]


def clone_template_slide(prs: Presentation, template_key: str):
    """Clone one of the boilerplate SFNL template slides into the active deck."""
    source = prs.slides[TEMPLATE_SLIDES[template_key] - 1]
    slide = prs.slides.add_slide(blank_layout(prs))

    src_bg = source._element.cSld.bg
    if src_bg is not None:
        cSld = slide._element.cSld
        if cSld.bg is not None:
            cSld.remove(cSld.bg)
        cSld.insert(0, deepcopy(src_bg))

    sp_tree = slide.shapes._spTree
    for shape in source.shapes:
        sp_tree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return slide


def clone_static_template_slide(prs: Presentation, template_key: str):
    """Clone only the non-placeholder artwork from a template slide."""
    source = prs.slides[TEMPLATE_SLIDES[template_key] - 1]
    slide = prs.slides.add_slide(blank_layout(prs))

    src_bg = source._element.cSld.bg
    if src_bg is not None:
        cSld = slide._element.cSld
        if cSld.bg is not None:
            cSld.remove(cSld.bg)
        cSld.insert(0, deepcopy(src_bg))

    sp_tree = slide.shapes._spTree
    for shape in source.shapes:
        if getattr(shape, "is_placeholder", False):
            continue
        sp_tree.insert_element_before(deepcopy(shape.element), "p:extLst")
    return slide, source


def find_placeholder(slide, idx: int):
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx == idx:
                return shape
        except Exception:
            continue
    raise KeyError(f"Placeholder idx {idx} niet gevonden op slide")


def find_shape(slide, name: str, occurrence: int = 0):
    matches = [shape for shape in slide.shapes if shape.name == name]
    if occurrence >= len(matches):
        raise KeyError(f"Shape {name!r} occurrence {occurrence} niet gevonden")
    return matches[occurrence]


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def _font_name(role: str) -> str:
    return STYLE["fonts"].get(role, STYLE["fonts"]["body"])


def apply_run_style(run, *, role="body", size=10, bold=False, italic=False, color=None) -> None:
    run.font.name = _font_name(role)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _clear_paragraph_bullets(paragraph) -> None:
    p_pr = paragraph._p.get_or_add_pPr()
    for child in list(p_pr):
        tag = child.tag.split("}")[-1]
        if tag.startswith("bu"):
            p_pr.remove(child)


def set_paragraphs(shape, paragraphs: list[dict], *, word_wrap=True) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = word_wrap

    if not paragraphs:
        return

    for idx, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        _clear_paragraph_bullets(p)
        p.level = spec.get("level", 0)
        if "alignment" in spec:
            p.alignment = spec["alignment"]
        run = p.add_run()
        run.text = spec.get("text", "")
        apply_run_style(
            run,
            role=spec.get("role", "body"),
            size=spec.get("size", 10),
            bold=spec.get("bold", False),
            italic=spec.get("italic", False),
            color=spec.get("color"),
        )


def set_text(shape, text: str, *, role="body", size=10, bold=False, italic=False, color=None) -> None:
    set_paragraphs(shape, [{
        "text": text,
        "role": role,
        "size": size,
        "bold": bold,
        "italic": italic,
        "color": color,
    }])


def replace_paragraph_text(paragraph, text: str) -> None:
    """Replace paragraph text while preserving existing paragraph/run formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.add_run().text = text


def set_text_preserve(shape, text: str) -> None:
    """Set a single-paragraph text frame without clearing its formatting."""
    tf = shape.text_frame
    replace_paragraph_text(tf.paragraphs[0], text)
    for paragraph in tf.paragraphs[1:]:
        replace_paragraph_text(paragraph, "")


def set_lines_preserve(shape, lines: list[str]) -> None:
    """Fill paragraph texts while keeping existing PowerPoint formatting intact."""
    tf = shape.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()

    for idx, line in enumerate(lines):
        replace_paragraph_text(tf.paragraphs[idx], line)

    for idx in range(len(lines), len(tf.paragraphs)):
        replace_paragraph_text(tf.paragraphs[idx], "")


def clear_run_highlights(shape) -> None:
    """Remove text highlight formatting from all runs in a shape."""
    if not getattr(shape, "has_text_frame", False):
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            r_pr = run._r.rPr
            if r_pr is None:
                continue
            for child in list(r_pr):
                if child.tag == qn("a:highlight"):
                    r_pr.remove(child)


def copy_table_cell_fill(source_cell, target_cell) -> None:
    """Copy only the fill definition from one table cell to another."""
    source_tc_pr = source_cell._tc.get_or_add_tcPr()
    target_tc_pr = target_cell._tc.get_or_add_tcPr()
    fill_tags = {
        "solidFill",
        "gradFill",
        "pattFill",
        "noFill",
        "blipFill",
        "grpFill",
    }

    for child in list(target_tc_pr):
        if child.tag.split("}")[-1] in fill_tags:
            target_tc_pr.remove(child)

    for child in source_tc_pr:
        if child.tag.split("}")[-1] in fill_tags:
            target_tc_pr.append(deepcopy(child))
