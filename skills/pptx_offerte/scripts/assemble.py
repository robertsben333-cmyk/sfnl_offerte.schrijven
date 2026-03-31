"""Assemble a PPTX from a slide_plan list and the SFNL base template."""
import os
import re
import shutil
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from skills.pptx_offerte.scripts.slides._utils import TEMPLATE_SLIDES

# Assets live alongside the consolidated skill package
_HERE = os.path.dirname(__file__)
ASSETS_DIR = os.path.normpath(os.path.join(_HERE, "../../pptx_offerte/assets"))
DEFAULT_BASE = os.path.join(ASSETS_DIR, "sfnl_base.pptx")

# Registry maps slide type → component add_slide function (loaded lazily)
_REGISTRY: dict | None = None

# Known types declared statically so we can validate without importing components
_KNOWN_TYPES = frozenset({
    "cover",
    "section_header",
    "aanleiding",
    "aanpak_overview",
    "fase_detail",
    "two_column",
    "tijdslijn",
    "team",
    "budget_table",
    "randvoorwaarden",
    "akkoord",
})


def _text_len(value) -> int:
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return len(str(value))
    if isinstance(value, str):
        return len(value.strip())
    if isinstance(value, list):
        return sum(_text_len(item) for item in value)
    if isinstance(value, dict):
        return sum(_text_len(item) for item in value.values())
    return len(str(value))


def _word_count(value) -> int:
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return len(str(value).split())
    if isinstance(value, str):
        return len(re.findall(r"\b[\w\-–]+\b", value))
    if isinstance(value, list):
        return sum(_word_count(item) for item in value)
    if isinstance(value, dict):
        return sum(_word_count(item) for item in value.values())
    return len(re.findall(r"\b[\w\-–]+\b", str(value)))


def _pick_size(value: int, tiers: list[tuple[int, int]], default: int) -> int:
    for limit, size in tiers:
        if value <= limit:
            return size
    return default


def _derive_style_context(slide_plan: list) -> dict:
    """Choose one consistent typography scale per slide group for the whole deck."""
    context: dict = {}

    fase_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "fase_detail"]
    if fase_slides:
        max_left = max(
            _word_count(content.get("acties_sfnl", [])[:3])
            + _word_count(content.get("acties_klant", [])[:3])
            + _word_count(content.get("deliverable", ""))
            + _word_count(content.get("tijdlijn", ""))
            + _word_count(content.get("dagen", ""))
            for content in fase_slides
        )
        max_right = max(
            _word_count(content.get("doel", "")) + _word_count(content.get("aanpak", ""))
            for content in fase_slides
        )
        context["fase_detail"] = {
            "left_label_size": _pick_size(max_left, [(35, 13), (50, 12), (70, 11)], 10),
            "left_body_size": _pick_size(max_left, [(35, 10), (50, 9), (70, 8)], 8),
            "right_label_size": _pick_size(max_right, [(130, 14), (170, 13), (220, 12)], 11),
            "right_body_size": _pick_size(max_right, [(130, 12), (170, 11), (220, 10)], 9),
        }

    overview_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "aanpak_overview"]
    if overview_slides:
        overview = overview_slides[0]
        max_desc = max((_word_count(phase.get("beschrijving", "")) for phase in overview.get("phases", [])), default=0)
        intro_len = _word_count(overview.get("intro", ""))
        context["aanpak_overview"] = {
            "intro_size": _pick_size(intro_len, [(28, 12), (45, 11), (70, 10)], 9),
            "desc_size": _pick_size(max_desc, [(22, 11), (35, 10), (55, 9)], 8),
            "timeline_size": 9,
        }

    aanleiding_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "aanleiding"]
    if aanleiding_slides:
        content = aanleiding_slides[0]
        max_block = max(
            _word_count(content.get("vraagstuk", "")),
            _word_count(content.get("uitdagingen", "")),
            _word_count(content.get("behoefte", "")),
        )
        context["aanleiding"] = {
            "body_size": _pick_size(max_block, [(45, 12), (65, 11), (85, 10)], 9),
            "summary_size": _pick_size(_word_count(content.get("summary_line", "")), [(18, 14), (28, 13), (40, 12)], 11),
        }

    randvoorwaarden_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "randvoorwaarden"]
    if randvoorwaarden_slides:
        total = max((_word_count(content.get("items", [])) for content in randvoorwaarden_slides), default=0)
        context["randvoorwaarden"] = {
            "body_size": _pick_size(total, [(55, 13), (80, 12), (110, 11)], 10),
        }

    budget_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "budget_table"]
    if budget_slides:
        total = max(
            (_word_count(content.get("tarief_note", "")) + _word_count(content.get("termijnen", [])) + _word_count(content.get("social_rate_disclaimer", "")) for content in budget_slides),
            default=0,
        )
        context["budget_table"] = {
            "notes_size": _pick_size(total, [(45, 10), (70, 9), (100, 8)], 8),
        }

    akkoord_slides = [entry.get("content", {}) for entry in slide_plan if entry.get("type") == "akkoord"]
    if akkoord_slides:
        total = max(
            (_word_count(content.get("randvoorwaarden_tekst", "")) + _word_count(content.get("termijnen", [])) for content in akkoord_slides),
            default=0,
        )
        context["akkoord"] = {
            "body_size": _pick_size(total, [(55, 12), (80, 11), (110, 10)], 9),
        }

    return context


def _load_registry() -> dict:
    global _REGISTRY
    if _REGISTRY is not None:
        return _REGISTRY
    from skills.pptx_offerte.scripts.slides import (
        cover,
        section_header,
        aanleiding,
        aanpak_overview,
        fase_detail,
        two_column,
        tijdslijn,
        team,
        budget_table,
        randvoorwaarden,
        akkoord,
    )
    _REGISTRY = {
        "cover": cover.add_slide,
        "section_header": section_header.add_slide,
        "aanleiding": aanleiding.add_slide,
        "aanpak_overview": aanpak_overview.add_slide,
        "fase_detail": fase_detail.add_slide,
        "two_column": two_column.add_slide,
        "tijdslijn": tijdslijn.add_slide,
        "team": team.add_slide,
        "budget_table": budget_table.add_slide,
        "randvoorwaarden": randvoorwaarden.add_slide,
        "akkoord": akkoord.add_slide,
    }
    return _REGISTRY


def assemble(slide_plan: list, output_path: str, base: str = DEFAULT_BASE) -> str:
    """
    Build a .pptx from slide_plan and save to output_path.

    slide_plan: [{"type": "cover", "content": {...}}, ...]
    Returns output_path.
    Raises ValueError for unknown slide types.
    """
    # Validate all types upfront before doing any heavy work.
    # Use the static set so unknown-type raises without importing components.
    for entry in slide_plan:
        slide_type = entry.get("type")
        if slide_type not in _KNOWN_TYPES:
            raise ValueError(
                f"Unknown slide type: {slide_type!r}. "
                f"Available: {sorted(_KNOWN_TYPES)}"
            )

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    if not slide_plan:
        prs = Presentation(base)
        prs.save(output_path)
        return output_path

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_dup, tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp_populated:
        duplicated_path = tmp_dup.name
        populated_path = tmp_populated.name

    try:
        boilerplate_count, duplicate_targets = _build_working_copy(base, slide_plan, duplicated_path)
        prs = Presentation(duplicated_path)
        registry = _load_registry()
        style_context = _derive_style_context(slide_plan)
        next_duplicate = 0
        for entry in slide_plan:
            content = dict(entry.get("content", {}))
            slide_type = entry["type"]
            content["__style_context__"] = style_context
            if slide_type in duplicate_targets:
                content["__target_slide_index__"] = boilerplate_count + next_duplicate
                next_duplicate += 1
            registry[slide_type](prs, content)
        prs.save(populated_path)
        _rewrite_slide_order(populated_path, output_path, boilerplate_count)
    finally:
        for temp_path in (duplicated_path, populated_path):
            if os.path.exists(temp_path):
                os.remove(temp_path)
    return output_path


def _build_working_copy(base: str, slide_plan: list, output_path: str) -> tuple[int, set[str]]:
    """Duplicate template slides into a working copy so we can edit the real slide contents."""
    duplicate_targets = {entry["type"] for entry in slide_plan if entry["type"] in TEMPLATE_SLIDES}
    boilerplate_count = len(Presentation(base).slides)
    if not duplicate_targets:
        shutil.copyfile(base, output_path)
        return boilerplate_count, duplicate_targets

    with tempfile.TemporaryDirectory() as unpack_dir:
        unpacked = Path(unpack_dir)
        with zipfile.ZipFile(base, "r") as src:
            src.extractall(unpacked)

        for entry in slide_plan:
            source_num = TEMPLATE_SLIDES.get(entry["type"])
            if source_num is not None:
                _duplicate_slide(unpacked, f"slide{source_num}.xml")

        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as dst:
            for file_path in sorted(unpacked.rglob("*")):
                if file_path.is_file():
                    dst.write(file_path, file_path.relative_to(unpacked).as_posix())

    return boilerplate_count, duplicate_targets


def _duplicate_slide(unpacked_dir: Path, source: str) -> None:
    slides_dir = unpacked_dir / "ppt" / "slides"
    rels_dir = slides_dir / "_rels"
    source_slide = slides_dir / source
    if not source_slide.exists():
        raise FileNotFoundError(source_slide)

    next_num = _get_next_slide_number(slides_dir)
    dest = f"slide{next_num}.xml"
    dest_slide = slides_dir / dest

    source_rels = rels_dir / f"{source}.rels"
    dest_rels = rels_dir / f"{dest}.rels"

    shutil.copy2(source_slide, dest_slide)
    if source_rels.exists():
        shutil.copy2(source_rels, dest_rels)
        rels_content = dest_rels.read_text(encoding="utf-8")
        rels_content = re.sub(
            r'\s*<Relationship[^>]*Type="[^"]*notesSlide"[^>]*/>\s*',
            "\n",
            rels_content,
        )
        dest_rels.write_text(rels_content, encoding="utf-8")

    _add_to_content_types(unpacked_dir, dest)
    rid = _add_to_presentation_rels(unpacked_dir, dest)
    slide_id = _get_next_slide_id(unpacked_dir)
    pres_path = unpacked_dir / "ppt" / "presentation.xml"
    pres_content = pres_path.read_text(encoding="utf-8")
    insertion = f'  <p:sldId id="{slide_id}" r:id="{rid}"/>\n</p:sldIdLst>'
    pres_content = pres_content.replace("</p:sldIdLst>", insertion)
    pres_path.write_text(pres_content, encoding="utf-8")


def _get_next_slide_number(slides_dir: Path) -> int:
    existing = [int(m.group(1)) for f in slides_dir.glob("slide*.xml") if (m := re.match(r"slide(\d+)\.xml", f.name))]
    return max(existing) + 1 if existing else 1


def _add_to_content_types(unpacked_dir: Path, dest: str) -> None:
    content_types_path = unpacked_dir / "[Content_Types].xml"
    content_types = content_types_path.read_text(encoding="utf-8")
    new_override = f'<Override PartName="/ppt/slides/{dest}" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
    if f"/ppt/slides/{dest}" not in content_types:
        content_types = content_types.replace("</Types>", f"  {new_override}\n</Types>")
        content_types_path.write_text(content_types, encoding="utf-8")


def _add_to_presentation_rels(unpacked_dir: Path, dest: str) -> str:
    pres_rels_path = unpacked_dir / "ppt" / "_rels" / "presentation.xml.rels"
    pres_rels = pres_rels_path.read_text(encoding="utf-8")
    rids = [int(m) for m in re.findall(r'Id="rId(\d+)"', pres_rels)]
    next_rid = max(rids) + 1 if rids else 1
    rid = f"rId{next_rid}"
    new_rel = f'<Relationship Id="{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/{dest}"/>'
    if f"slides/{dest}" not in pres_rels:
        pres_rels = pres_rels.replace("</Relationships>", f"  {new_rel}\n</Relationships>")
        pres_rels_path.write_text(pres_rels, encoding="utf-8")
    return rid


def _get_next_slide_id(unpacked_dir: Path) -> int:
    pres_path = unpacked_dir / "ppt" / "presentation.xml"
    pres_content = pres_path.read_text(encoding="utf-8")
    slide_ids = [int(m) for m in re.findall(r'<p:sldId[^>]*id="(\d+)"', pres_content)]
    return max(slide_ids) + 1 if slide_ids else 256


def _rewrite_slide_order(source_path: str, output_path: str, boilerplate_count: int) -> None:
    """Write a reordered copy of the pptx without mutating slide parts in-memory."""
    presentation_xml = "ppt/presentation.xml"
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}

    with zipfile.ZipFile(source_path, "r") as src:
        xml_bytes = src.read(presentation_xml)
        root = ET.fromstring(xml_bytes)
        sld_id_list = root.find("p:sldIdLst", ns)
        if sld_id_list is None:
            raise ValueError("presentation.xml mist p:sldIdLst")

        slide_ids = list(sld_id_list)
        generated_slide_ids = slide_ids[boilerplate_count:]
        sld_id_list[:] = generated_slide_ids
        updated_xml = ET.tostring(root, encoding="utf-8", xml_declaration=True)

        with zipfile.ZipFile(output_path, "w", zipfile.ZIP_DEFLATED) as dst:
            for item in src.infolist():
                if item.filename == presentation_xml:
                    dst.writestr(item, updated_xml)
                else:
                    dst.writestr(item, src.read(item.filename))
