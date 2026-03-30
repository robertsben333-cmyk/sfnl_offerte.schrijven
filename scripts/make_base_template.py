"""Strip MBC template down to slide master + boilerplate slides (17+)."""
import os
from pptx import Presentation

SRC = "templates/offerte_mbc_template.pptx"
DST = "skills/pptx-offerte/assets/sfnl_base.pptx"

import sys
if not os.path.exists(SRC):
    print(f"ERROR: source template not found: {SRC}", file=sys.stderr)
    print("Run this script from the project root directory.", file=sys.stderr)
    sys.exit(2)

prs = Presentation(SRC)

# We need to remove the first 16 slides (indices 0-15).
# python-pptx doesn't have a remove_slide method, so we manipulate the XML directly.
# The slide list is in prs.slides._sldIdLst
xml_slides = prs.slides._sldIdLst
slide_ids = list(xml_slides)

print(f"Total slides in source: {len(slide_ids)}")

# Note: This removes the slide from the presentation index but leaves orphaned
# XML parts in the zip (known python-pptx limitation). The output is a valid
# PPTX; PowerPoint ignores orphaned parts.
# Remove slides at indices 0-15 (first 16 slides)
for sldId in slide_ids[:16]:
    xml_slides.remove(sldId)

os.makedirs(os.path.dirname(DST), exist_ok=True)
prs.save(DST)
print(f"Base template saved to {DST}")
print(f"Slides remaining: {len(prs.slides)}")
