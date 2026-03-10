#!/usr/bin/env python3
"""
SFNL Offerte Generator
Generates a proposal PPTX from the standard template using a JSON config file.

Usage:
  py generate_offerte.py <config.json> <output.pptx>
"""

import copy
import json
import sys
from pathlib import Path

import lxml.etree as etree
from pptx import Presentation
from pptx.oxml.ns import nsmap

TEMPLATE_PATH = Path(__file__).parent.parent / "templates" / "offerte_mbc_template.pptx"

# Slide indices (0-based) — valid BEFORE any extra slides are inserted
SLIDE_COVER         = 0
SLIDE_TOC           = 1
SLIDE_AANLEIDING    = 2
SLIDE_AANPAK_HDR    = 3
SLIDE_AANPAK_OVW    = 5   # "ONZE AANPAK" overview with 3 chevrons
SLIDE_FASE1         = 6
SLIDE_FASE2         = 7
SLIDE_FASE3         = 8
SLIDE_TIJDSLIJN     = 9   # shifts right when extra fase slides are inserted
SLIDE_TEAM_HDR      = 11
SLIDE_TEAM          = 12
SLIDE_BEGROTING_HDR = 13
SLIDE_BEGROTING     = 14
SLIDE_RANDVW        = 15
SLIDE_SFNL_START    = 16  # All slides from here on are fixed SFNL boilerplate

NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


# ─── Formatting helpers ───────────────────────────────────────────────────────

def format_eur(amount: float) -> str:
    """Format a euro amount in Dutch style: € 41.440 (dot as thousands separator)."""
    formatted = f"{amount:,.0f}".replace(",", ".")
    return f"€ {formatted}"


def _number_to_dutch_upper(n: int) -> str:
    mapping = {1: "ÉÉN", 2: "TWEE", 3: "DRIE", 4: "VIER", 5: "VIJF",
               6: "ZES", 7: "ZEVEN", 8: "ACHT"}
    return mapping.get(n, str(n))


# ─── Text helpers ─────────────────────────────────────────────────────────────

def para_full_text(para) -> str:
    return "".join(r.text for r in para.runs)


def set_para_text(para, new_text: str):
    """Set paragraph text in first run, clear all other runs."""
    if not para.runs:
        para.add_run().text = new_text
        return
    para.runs[0].text = new_text
    for run in para.runs[1:]:
        run.text = ""


def replace_in_tf(tf, old: str, new: str):
    """Replace a string across all paragraphs in a text frame, handling split runs."""
    for para in tf.paragraphs:
        full = para_full_text(para)
        if old in full:
            set_para_text(para, full.replace(old, new))


def replace_in_slide(slide, replacements: dict):
    """Apply multiple replacements across all text frames and tables on a slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for old, new in replacements.items():
                replace_in_tf(shape.text_frame, old, new)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for old, new in replacements.items():
                        replace_in_tf(cell.text_frame, old, new)


def set_shape_full_text(shape, lines: list):
    """
    Overwrite a shape's text frame with the given lines (one per paragraph).
    Preserves run formatting of the first template paragraph for all new paragraphs.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    txBody = tf._txBody

    existing = txBody.findall(f"{{{NS}}}p")
    template_p = copy.deepcopy(existing[0]) if existing else None

    for p in existing:
        txBody.remove(p)

    template_runs = template_p.findall(f"{{{NS}}}r") if template_p is not None else []

    for line in lines:
        new_p = copy.deepcopy(template_p) if template_p is not None else etree.SubElement(txBody, f"{{{NS}}}p")

        for r in new_p.findall(f"{{{NS}}}r"):
            new_p.remove(r)

        if template_runs:
            new_r = copy.deepcopy(template_runs[0])
            t_el = new_r.find(f"{{{NS}}}t")
            if t_el is None:
                t_el = etree.SubElement(new_r, f"{{{NS}}}t")
            t_el.text = line
        else:
            new_r = etree.SubElement(new_p, f"{{{NS}}}r")
            t_el = etree.SubElement(new_r, f"{{{NS}}}t")
            t_el.text = line

        new_p.append(new_r)
        txBody.append(new_p)


# ─── Slide duplication ────────────────────────────────────────────────────────

def duplicate_and_insert_slide(prs, source_index: int, insert_before_index: int):
    """
    Clone the slide at source_index and insert the clone before insert_before_index.
    All shapes and text are preserved. Returns the new Slide object.

    Note: Image/media relationship references are copied at the XML level but the
    underlying Part relationships are not duplicated. Text-only slides (e.g. fase slides)
    work correctly; slides with embedded images may show broken references.
    """
    src_slide = prs.slides[source_index]

    # Add a new slide at the end using the same layout
    new_slide = prs.slides.add_slide(src_slide.slide_layout)

    # Replace new slide's shape tree with a deep copy of the source's shape tree
    src_sp_tree = src_slide._element.spTree
    dst_sp_tree = new_slide._element.spTree
    for child in list(dst_sp_tree):
        dst_sp_tree.remove(child)
    for child in src_sp_tree:
        dst_sp_tree.append(copy.deepcopy(child))

    # Reorder _sldIdLst: move the new slide (currently last) to insert_before_index
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    new_sldId = sldIds[-1]

    if insert_before_index < len(sldIds) - 1:
        sldIds[insert_before_index].addprevious(new_sldId)
    # else: already at the end — no reorder needed

    return new_slide


# ─── Config validation ────────────────────────────────────────────────────────

def validate_config(cfg: dict):
    """Raise ValueError with a clear message if required fields are missing."""
    errors = []

    def require(path, obj=cfg):
        keys = path.split(".")
        cur = obj
        for k in keys:
            if not isinstance(cur, dict) or k not in cur:
                errors.append(f"Missing required field: {path}")
                return
            cur = cur[k]

    require("client_name")
    require("proposal_date")
    require("aanleiding.summary_line")
    require("aanleiding.maatschappelijk_vraagstuk")
    require("aanleiding.behoefte_van_klant")
    require("aanpak.fases")
    require("tijdslijn.maanden")
    require("begroting.total_excl_btw")
    require("begroting.betaaltermijnen")

    fases = cfg.get("aanpak", {}).get("fases", [])
    if not fases:
        errors.append("aanpak.fases must contain at least one fase")
    for i, fase in enumerate(fases):
        for field in ("number", "name", "dagen", "deliverable"):
            if field not in fase:
                errors.append(f"aanpak.fases[{i}] missing field: {field}")

    if errors:
        raise ValueError("Config validation failed:\n" + "\n".join(f"  - {e}" for e in errors))


# ─── Slide modifiers ──────────────────────────────────────────────────────────

def modify_aanleiding(slide, cfg):
    """Slide 3: AANLEIDING OFFERTE — summary line + 3 content blocks."""
    a = cfg["aanleiding"]
    replace_in_slide(slide, {"[SAMENVATTENDE ZIN]": a["summary_line"]})

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = shape.text_frame.text
        if "Maatschappelijk vraagstuk" in txt:
            set_shape_full_text(shape, ["Maatschappelijk vraagstuk", a["maatschappelijk_vraagstuk"]])
        elif "Grootste uitdagingen" in txt:
            set_shape_full_text(shape, ["Grootste uitdagingen", a.get("grootste_uitdagingen", "")])
        elif "Behoefte van de klant" in txt:
            set_shape_full_text(shape, ["Behoefte van de klant\xa0", a["behoefte_van_klant"]])


def modify_aanpak_overview(slide, cfg):
    """Slide 6: ONZE AANPAK — subtitle, chevron labels, phase descriptions, timeline dates."""
    fases = cfg["aanpak"]["fases"]
    client = cfg["client_name"]

    # Subtitle
    subtitle = cfg["aanpak"].get(
        "overview_subtitle",
        f"IN {_number_to_dutch_upper(len(fases))} FASES BRENGEN WE DE IMPACT VAN "
        f"{client.upper()} IN KAART EN LEGGEN WE DE BASIS VOOR EEN DUURZAAM FINANCIERINGSMODEL"
    )
    for shape in slide.shapes:
        if shape.has_text_frame and "IN DRIE FASES" in shape.text_frame.text:
            set_shape_full_text(shape, [subtitle])

    # Chevron labels — matched by current "FASE N:" prefix
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        for i, fase in enumerate(fases[:3]):
            if t.startswith(f"FASE {i + 1}:"):
                set_shape_full_text(shape, [f"FASE {fase['number']}: {fase['name'].upper()}"])
                break

    # Phase description rectangles — matched by position order among Rectangle shapes
    # The template has exactly 3 Rectangle shapes with descriptions, in phase order
    desc_rects = [
        s for s in slide.shapes
        if "Rectangle" in s.name and s.has_text_frame
        and len(s.text_frame.text) > 30
    ]
    for i, shape in enumerate(desc_rects[:len(fases)]):
        desc = fases[i].get("overview_description")
        if desc:
            set_shape_full_text(shape, [desc])

    # Timeline date labels — TextBox/Tekstvak shapes that contain digits
    date_shapes = [
        s for s in slide.shapes
        if s.has_text_frame
        and ("Tekstvak" in s.name or "TextBox" in s.name)
        and any(c.isdigit() for c in s.text_frame.text)
    ]
    for i, shape in enumerate(date_shapes[:len(fases)]):
        if fases[i].get("tijdlijn"):
            set_shape_full_text(shape, [fases[i]["tijdlijn"]])


def _modify_fase_slide(slide, fase: dict, client_name: str):
    """Modify a single fase detail slide."""
    n = fase["number"]
    name = fase["name"].upper()
    dagen = fase["dagen"]
    doel = fase.get("doel", "")
    acties_sfnl = fase.get("acties_sfnl", [])
    acties_klant = fase.get("acties_klant", [])
    deliverable = fase.get("deliverable", "")
    outcomes_note = fase.get("outcomes_note", "")

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text

        if shape.name == "Title 4":
            set_shape_full_text(shape, [f"{n}. {name}"])

        elif t.strip().startswith("Doel"):
            lines = ["Doel", doel]
            if fase.get("aanpak"):
                lines += ["", "Aanpak", fase["aanpak"]]
            set_shape_full_text(shape, lines)

        elif "Acties Social Finance NL" in t:
            lines = ["Acties Social Finance NL:"]
            lines += [f"   {a}" for a in acties_sfnl]
            if outcomes_note:
                lines += ["", outcomes_note]
            lines += ["", f"Acties {client_name}:"]
            lines += [f"   {a}" for a in acties_klant]
            lines += ["", f"Deliverable fase {n}: {deliverable}", "", f"Duur: {dagen} dagen"]
            set_shape_full_text(shape, lines)

        elif t.strip().startswith("LET OP:") or t.strip().startswith("LET OP "):
            set_shape_full_text(shape, [""])


def insert_extra_fase(prs, fase: dict, client_name: str, extra_count: int = 0):
    """
    Insert a single extra fase slide (fase 4+) by cloning the fase 3 template slide
    and placing it before the tijdslijn slide.

    Call this manually for each extra fase after running generate(), or use it from
    a wrapper script. Pass extra_count = number of extra slides already inserted so
    the insertion index stays correct.

    Returns the new Slide object.
    """
    insert_before = SLIDE_TIJDSLIJN + extra_count
    new_slide = duplicate_and_insert_slide(prs, SLIDE_FASE3, insert_before)
    _modify_fase_slide(new_slide, fase, client_name)
    return new_slide


def modify_fase_slides(slides, cfg):
    """Modify fase 1–3 detail slides at their fixed indices."""
    fases = cfg["aanpak"]["fases"]
    fixed_indices = [SLIDE_FASE1, SLIDE_FASE2, SLIDE_FASE3]
    for i, fase in enumerate(fases[:3]):
        _modify_fase_slide(slides[fixed_indices[i]], fase, cfg["client_name"])


def modify_tijdslijn(slide, cfg):
    """Slide 10: TIJDSLIJN — replace header text."""
    maanden = cfg["tijdslijn"]["maanden"]
    header = cfg["tijdslijn"].get(
        "header",
        f"BINNEN {_number_to_dutch_upper(maanden)} MAANDEN STELLEN WE DE MAATSCHAPPELIJKE "
        f"BUSINESSCASE OP EN BRENGEN WE DE HAALBAARHEID VAN EEN DUURZAAM FINANCIERINGSMODEL IN KAART"
    )
    for shape in slide.shapes:
        if shape.has_text_frame and "BINNEN" in shape.text_frame.text and "MAANDEN" in shape.text_frame.text:
            set_shape_full_text(shape, [header])


def modify_team(slide, cfg):
    """Slide 13: TEAM — replace 3 team member blocks."""
    members = cfg.get("team", [])

    name_shapes = [s for s in slide.shapes if s.has_text_frame and "[TEAMLID]" in s.text_frame.text]
    bio_shapes  = [s for s in slide.shapes if s.has_text_frame and "[Cv-omschrijving]" in s.text_frame.text]

    for i, member in enumerate(members[:3]):
        if i < len(name_shapes):
            set_shape_full_text(name_shapes[i], [member["name"].upper(), member["title_short"]])
        if i < len(bio_shapes):
            set_shape_full_text(bio_shapes[i], [member["bio"]])


def modify_begroting(slide, cfg):
    """Slide 15: BEGROTING — budget line items and totals."""
    b = cfg["begroting"]
    rows = b.get("rows", [])
    total = b["total_excl_btw"]
    btw_pct = b.get("btw_percentage", 21)
    btw_amount = b.get("btw", round(total * btw_pct / 100))
    total_incl = b.get("total_incl_btw", total + btw_amount)

    budget_lines = [
        "Social Finance NL voert de opdracht uit op basis van onze Algemene Voorwaarden.",
        "Het tarief is exclusief BTW en op basis van 8 uur per dag.",
        "Dit tarief is inclusief reiskosten binnen Nederland.",
        "",
    ]
    for row in rows:
        budget_lines.append(
            f"{row['fase']}: {row['dagen']} dagen x {format_eur(row['tarief'])} = {format_eur(row['totaal'])}"
        )
    budget_lines += [
        "",
        f"Totaal (excl. BTW): {format_eur(total)}",
        f"BTW ({btw_pct}%): {format_eur(btw_amount)}",
        f"Totaal (incl. BTW): {format_eur(total_incl)}",
    ]

    for shape in slide.shapes:
        if shape.has_text_frame and "Algemene Voorwaarden" in shape.text_frame.text:
            set_shape_full_text(shape, budget_lines)
            break


def modify_randvoorwaarden(slide, cfg):
    """Slide 16: RANDVOORWAARDEN EN AKKOORD — payment schedule + client name."""
    b = cfg["begroting"]
    client = cfg["client_name"]

    replace_in_slide(slide, {"[naam klant]": client})

    termijn_lines = [
        f"{format_eur(t['amount'])} — {t['description']}"
        for t in b.get("betaaltermijnen", [])
    ]

    for shape in slide.shapes:
        if shape.has_text_frame and "Randvoorwaarden" in shape.text_frame.text:
            lines = [
                "Randvoorwaarden",
                "Zodra er akkoord is op de offerte wordt de definitieve teamsamenstelling vastgesteld.",
                f"De planning is indicatief en wordt vastgesteld in overleg met {client} na akkoord op de offerte.",
                "Zonder vooraf verkregen toestemming zal SFNL geen meerwerk in rekening brengen.",
                "",
                "Facturatieschema:",
            ] + termijn_lines
            if cfg.get("factuuradres"):
                lines += ["", f"Factuuradres: {cfg['factuuradres']}"]
            set_shape_full_text(shape, lines)
            break


# ─── Main ─────────────────────────────────────────────────────────────────────

def generate(config_path: str, output_path: str):
    with open(config_path, encoding="utf-8") as f:
        cfg = json.load(f)

    validate_config(cfg)

    prs = Presentation(str(TEMPLATE_PATH))

    fases = cfg["aanpak"]["fases"]

    print(f"Template loaded: {len(prs.slides)} slides")
    print(f"Client:  {cfg['client_name']}")
    print(f"Fases:   {len(fases)}")
    print(f"Tarief:  {format_eur(cfg.get('day_rate', 1480))}/dag")

    # Global replacements across all editable slides (fixed indices — no extra slides yet)
    global_replacements = {
        "[naam klant]":   cfg["client_name"],
        "[naam project]": cfg["client_name"],
        "[NAAM PROJECT]": cfg["client_name"].upper(),
    }
    for i in range(SLIDE_SFNL_START):
        replace_in_slide(prs.slides[i], global_replacements)

    modify_aanleiding(prs.slides[SLIDE_AANLEIDING], cfg)
    modify_aanpak_overview(prs.slides[SLIDE_AANPAK_OVW], cfg)
    modify_fase_slides(prs.slides, cfg)
    modify_tijdslijn(prs.slides[SLIDE_TIJDSLIJN], cfg)
    modify_team(prs.slides[SLIDE_TEAM], cfg)
    modify_begroting(prs.slides[SLIDE_BEGROTING], cfg)
    modify_randvoorwaarden(prs.slides[SLIDE_RANDVW], cfg)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    b = cfg["begroting"]
    print(f"\nSaved:   {out}")
    print(f"Budget:  {format_eur(b['total_excl_btw'])} excl. BTW  |  {format_eur(b.get('total_incl_btw', b['total_excl_btw']))} incl. BTW")
    print(f"Slides:  {len(prs.slides)} total")

    if len(fases) > 3:
        extra = fases[3:]
        print(f"\n⚠️  {len(extra)} extra fase(s) in config — NOT auto-inserted.")
        print("   The overview slide (slide 6) only has 3 chevrons and needs manual adjustment.")
        print("   To insert each extra fase slide, use insert_extra_fase() from this module,")
        print("   or ask the agent to handle it after reviewing the output.")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: py generate_offerte.py <config.json> <output.pptx>")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2])
