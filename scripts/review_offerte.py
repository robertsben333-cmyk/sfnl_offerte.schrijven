#!/usr/bin/env python3
"""
SFNL Offerte Reviewer
Checks a generated PPTX for budget consistency, placeholder text, and writing quality.

Usage:
  py review_offerte.py <output.pptx> [config.json]

Exit codes:
  0 — No issues found
  1 — Issues found (see report)
  2 — File error (cannot open PPTX)
"""

import json
import re
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation


FORBIDDEN = [
    "alsmede", "tevens", "waarbij", "derhalve", "teneinde",
    "ten behoeve van", "in het kader van", "op het gebied van",
    "gedegen", "robuust", "uitstekend", "innovatief",
    "wij zijn verheugd", "graag stellen wij voor", "hierbij",
]

PLACEHOLDERS = [
    "[naam klant]", "[TEAMLID]", "[Cv-omschrijving]",
    "[SAMENVATTENDE ZIN]", "[naam project]", "[NAAM PROJECT]",
]

# Slides with fixed SFNL boilerplate — skip for content checks
SFNL_BOILERPLATE_START = 16  # slides 17+ are fixed


# ─── Text extraction ──────────────────────────────────────────────────────────

def _slide_text(slide) -> str:
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    parts.append(cell.text_frame.text)
    return "\n".join(parts)


def extract_slides(pptx_path: str):
    """Returns list of (slide_number_1based, text) for editable slides."""
    prs = Presentation(pptx_path)
    results = []
    limit = min(SFNL_BOILERPLATE_START, len(prs.slides))
    for i in range(limit):
        results.append((i + 1, _slide_text(prs.slides[i])))
    return results


# ─── Budget checks ────────────────────────────────────────────────────────────

def check_budget(cfg) -> list:
    issues = []

    if isinstance(cfg, list):
        budget_entry = next((entry for entry in cfg if entry.get("type") == "budget_table"), None)
        if budget_entry is None:
            return issues

        rows = budget_entry.get("content", {}).get("rows", [])
        total_row = None
        subtotal = 0
        for row in rows:
            fase = str(row.get("fase", "")).lower()
            kosten = row.get("kosten")
            if not isinstance(kosten, (int, float)):
                continue
            if "totaal" in fase:
                total_row = kosten
            else:
                subtotal += kosten
        if total_row is not None and abs(subtotal - total_row) > 1:
            issues.append(f"Totaalrij ({total_row}) ≠ som van niet-totaalrijen ({subtotal})")
        return issues

    b = cfg.get("begroting", {})
    rows = b.get("rows", [])

    computed_total = 0
    for row in rows:
        expected = row["dagen"] * row["tarief"]
        computed_total += expected
        actual = row.get("totaal", 0)
        if abs(actual - expected) > 0:
            issues.append(
                f"Rij '{row['fase']}': {row['dagen']}d × {row['tarief']} = {expected}, "
                f"maar totaal staat op {actual}"
            )

    declared_total = b.get("total_excl_btw", 0)
    if abs(computed_total - declared_total) > 1:
        issues.append(
            f"total_excl_btw ({declared_total}) ≠ som van rijen ({computed_total})"
        )

    btw_pct = b.get("btw_percentage", 21)
    expected_btw = round(declared_total * btw_pct / 100)
    actual_btw = b.get("btw", 0)
    if abs(actual_btw - expected_btw) > 1:
        issues.append(
            f"BTW ({actual_btw}) ≠ {btw_pct}% × {declared_total} = {expected_btw}"
        )

    expected_incl = declared_total + actual_btw
    actual_incl = b.get("total_incl_btw", 0)
    if abs(actual_incl - expected_incl) > 1:
        issues.append(
            f"total_incl_btw ({actual_incl}) ≠ {declared_total} + {actual_btw} = {expected_incl}"
        )

    sum_termijnen = sum(t.get("amount", 0) for t in b.get("betaaltermijnen", []))
    if actual_incl and abs(sum_termijnen - actual_incl) > 2:
        issues.append(
            f"Betaaltermijnen som ({sum_termijnen}) ≠ total_incl_btw ({actual_incl})"
        )

    return issues


# ─── Content checks ───────────────────────────────────────────────────────────

def check_placeholders(slides: list) -> list:
    issues = []
    for slide_num, text in slides:
        for ph in PLACEHOLDERS:
            if ph.lower() in text.lower():
                issues.append(f"Slide {slide_num}: placeholder '{ph}' nog aanwezig")
    return issues


def check_forbidden_words(slides: list) -> list:
    issues = []
    seen = set()
    for slide_num, text in slides:
        text_lower = text.lower()
        for word in FORBIDDEN:
            key = (slide_num, word)
            if key not in seen and word in text_lower:
                issues.append(f"Slide {slide_num}: verboden woord/zin '{word}'")
                seen.add(key)
    return issues


def check_long_sentences(slides: list, max_words: int = 35) -> list:
    issues = []
    for slide_num, text in slides:
        for sentence in re.split(r"[.!?]\s+", text):
            sentence = sentence.strip()
            if not sentence:
                continue
            word_count = len(sentence.split())
            if word_count > max_words:
                preview = sentence[:80] + ("…" if len(sentence) > 80 else "")
                issues.append(
                    f"Slide {slide_num}: zin heeft {word_count} woorden: \"{preview}\""
                )
    return issues


def check_em_dashes(slides: list) -> list:
    """Flag em-dashes (–) used as bullet-point starters."""
    issues = []
    for slide_num, text in slides:
        for line in text.splitlines():
            line = line.strip()
            if line.startswith("–") or line.startswith("—"):
                preview = line[:60] + ("…" if len(line) > 60 else "")
                issues.append(f"Slide {slide_num}: em-dash als opsomteken: \"{preview}\"")
    return issues


# ─── Report ───────────────────────────────────────────────────────────────────

def _section(title: str, issues: list, width: int = 60):
    label = f"[{title}]"
    if issues:
        print(f"\n{label}")
        for issue in issues:
            print(f"  ✗ {issue}")
    else:
        print(f"\n{label} {'—':>{width - len(label) - 2}} ✓")


def review(pptx_path: str, config_path: Optional[str] = None):
    pptx_file = Path(pptx_path)
    if not pptx_file.exists():
        print(f"FOUT: bestand niet gevonden: {pptx_path}")
        sys.exit(2)

    try:
        slides = extract_slides(pptx_path)
    except Exception as exc:
        print(f"FOUT: kan PPTX niet openen: {exc}")
        sys.exit(2)

    width = 60
    print(f"\n{'=' * width}")
    print(f"SFNL OFFERTE REVIEW — {pptx_file.name}")
    print(f"{'=' * width}")
    print(f"Gecontroleerde slides: {len(slides)} (boilerplate overgeslagen)")

    results = {}

    # Budget (only if config provided)
    if config_path:
        try:
            with open(config_path, encoding="utf-8") as f:
                cfg = json.load(f)
            results["BEGROTING"] = check_budget(cfg)
        except Exception as exc:
            print(f"\nWAARSCHUWING: config niet gelezen: {exc}")
            results["BEGROTING"] = [f"Config niet beschikbaar: {exc}"]
    else:
        results["BEGROTING"] = ["Geen config opgegeven — budgetcontrole overgeslagen"]

    results["PLACEHOLDERS"]    = check_placeholders(slides)
    results["VERBODEN WOORDEN"] = check_forbidden_words(slides)
    results["LANGE ZINNEN"]     = check_long_sentences(slides)
    results["EM-DASHES"]        = check_em_dashes(slides)

    for title, issues in results.items():
        _section(title, issues, width)

    total_errors = sum(
        len(v) for k, v in results.items()
        if k != "BEGROTING" or (config_path and v)
    )
    # Budget "no config" is informational, not an error
    if not config_path:
        total_errors -= len(results.get("BEGROTING", []))

    print(f"\n{'=' * width}")
    if total_errors == 0:
        print("RESULTAAT: Alle controles geslaagd ✓")
        sys.exit(0)
    else:
        print(f"RESULTAAT: {total_errors} aandachtspunt(en) gevonden")
        sys.exit(1)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Gebruik: py review_offerte.py <output.pptx> [config.json]")
        sys.exit(2)
    review(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else None)
