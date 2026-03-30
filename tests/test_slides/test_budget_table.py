"""Tests for PPTX budget_table slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"

CONTENT = {
    "rows": [
        {"fase": "Fase 1: Testfase", "dagen": 8, "kosten": 11840},
        {"fase": "Fase 2: Testfase", "dagen": 16, "kosten": 23680},
        {"fase": "Totaal", "dagen": 24, "kosten": 35520},
    ],
    "day_rate": 1480,
    "tarief_note": "Voor deze opdracht hanteren we een gereduceerd tarief van €1.200 per dag.",
    "social_rate_disclaimer": "Dit tarief is niet geldig voor vervolgopdrachten of andere opdrachten.",
    "termijnen": ["50% bij opdrachtverlening", "50% bij oplevering"],
    "proposition": "mbc"
}


def test_budget_table_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.budget_table import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_budget_table_has_table_shape():
    from skills.pptx_offerte.scripts.slides.budget_table import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    slide = prs.slides[-1]
    assert any(s.has_table for s in slide.shapes)


def test_budget_table_contains_fase_names():
    from skills.pptx_offerte.scripts.slides.budget_table import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    slide = prs.slides[-1]
    table_shape = next(s for s in slide.shapes if s.has_table)
    all_cells = [cell.text for row in table_shape.table.rows for cell in row.cells]
    assert any("Fase 1" in c for c in all_cells)
    assert any("Fase 2" in c for c in all_cells)


def test_budget_table_contains_tarief_note_and_disclaimer():
    from skills.pptx_offerte.scripts.slides.budget_table import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = " ".join(s.text_frame.text for s in prs.slides[-1].shapes if s.has_text_frame)
    assert "gereduceerd tarief" in text
    assert "niet geldig voor vervolgopdrachten" in text
