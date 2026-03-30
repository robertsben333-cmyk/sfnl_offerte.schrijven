"""Tests for PPTX randvoorwaarden slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"

CONTENT = {
    "title": "RANDVOORWAARDEN",
    "items": ["Eerste randvoorwaarde.", "Tweede randvoorwaarde."],
    "proposition": "mbc"
}


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_randvoorwaarden_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.randvoorwaarden import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_randvoorwaarden_contains_items():
    from skills.pptx_offerte.scripts.slides.randvoorwaarden import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Eerste randvoorwaarde" in text
    assert "Tweede randvoorwaarde" in text
