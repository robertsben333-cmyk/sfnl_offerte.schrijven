"""Tests for PPTX fase_detail slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"

CONTENT = {
    "number": 1,
    "naam": "Testfase",
    "klant": "Testklant",
    "doel": "Het doel is om X te bereiken.",
    "aanpak": "De aanpak bestaat uit Y stappen.",
    "acties_sfnl": ["Deskresearch", "Stakeholderinterviews"],
    "acties_klant": ["Beschikbaar stellen gegevens"],
    "deliverable": "Rapport en presentatie",
    "dagen": 8,
    "tijdlijn": "jan–feb 2026",
    "proposition": "mbc"
}


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_fase_detail_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_fase_detail_contains_phase_name():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "Testfase" in _all_text(prs.slides[-1])


def test_fase_detail_contains_doel_and_aanpak():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "doel is om X" in text
    assert "aanpak bestaat" in text


def test_fase_detail_contains_acties():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Deskresearch" in text
    assert "Stakeholderinterviews" in text
    assert "TESTKLANT" in text


def test_fase_detail_contains_deliverable_and_days():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Rapport en presentatie" in text
    assert "8" in text


def test_fase_detail_labels_are_bold():
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide

    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    slide = prs.slides[-1]
    right_box = next(ph for ph in slide.placeholders if ph.placeholder_format.idx == 11)
    left_box = next(ph for ph in slide.placeholders if ph.placeholder_format.idx == 12)

    right_labels = {
        paragraph.text.strip(): paragraph.runs[0].font.bold
        for paragraph in right_box.text_frame.paragraphs
        if paragraph.text.strip() in {"Doel", "Aanpak"}
    }
    left_labels = {
        paragraph.text.strip(): paragraph.runs[0].font.bold
        for paragraph in left_box.text_frame.paragraphs
        if paragraph.text.strip() in {"Acties Social Finance NL", "Acties TESTKLANT", "Deliverable", "Duur"}
    }

    assert right_labels == {"Doel": True, "Aanpak": True}
    assert left_labels == {
        "Acties Social Finance NL": True,
        "Acties TESTKLANT": True,
        "Deliverable": True,
        "Duur": True,
    }


def test_fase_detail_accepts_dotted_substep_number():
    """Substep slides use dotted numbers like '1.1' — title must render as '1.1. Naam'."""
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    content = {
        "number": "1.1",
        "naam": "VERANDERTHEORIE EN EFFECTENKAART",
        "subtitle": "FASE 1: IMPACTNARRATIEF",
        "klant": "Testklant",
        "doel": "Verandertheorie en effectenkaart opstellen.",
        "aanpak": "Kick-off, stakeholderinterviews en werksessie.",
        "acties_sfnl": ["Verandertheorie opstellen", "Effectenkaart bouwen"],
        "acties_klant": ["Deelnemen aan werksessie"],
        "deliverable": "Validated effectenkaart",
        "dagen": 4,
        "tijdlijn": "jan 2026",
        "proposition": "mbc",
    }
    add_slide(prs, content)
    slide = prs.slides[-1]
    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "1.1" in text
    assert "VERANDERTHEORIE EN EFFECTENKAART" in text
    assert "FASE 1: IMPACTNARRATIEF" in text


def test_fase_detail_substep_subtitle_overrides_default():
    """When subtitle is provided, it replaces the default 'AANPAK FASE X' label."""
    from skills.pptx_offerte.scripts.slides.fase_detail import add_slide
    prs = Presentation(BASE)
    content = {
        "number": "2.2",
        "naam": "BUSINESSCASE OPSTELLEN EN VALIDEREN",
        "subtitle": "FASE 2: MAATSCHAPPELIJKE BUSINESSCASE",
        "klant": "Testklant",
        "doel": "Businesscase opstellen.",
        "aanpak": "Data waarderen en rapport schrijven.",
        "acties_sfnl": ["Financiële waardering"],
        "acties_klant": ["Reviewronde rapport"],
        "deliverable": "Definitieve businesscase",
        "dagen": 9,
        "tijdlijn": "apr–mei 2026",
        "proposition": "mbc",
    }
    add_slide(prs, content)
    slide = prs.slides[-1]
    text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "FASE 2: MAATSCHAPPELIJKE BUSINESSCASE" in text
    assert "AANPAK FASE 2.2" not in text
