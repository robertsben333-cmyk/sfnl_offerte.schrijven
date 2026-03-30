"""Tests for PPTX team slide component."""
import os, sys, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"
CONTENT = {
    "proposition": "mbc",
    "members": [
        {"name": "Laura Brouwer", "title": "ASSOCIATE DIRECTOR", "bio": "Laura heeft 10 jaar ervaring."},
        {"name": "Dieuwertje Roos", "title": "MANAGER", "bio": "Dieuwertje is gespecialiseerd in MBC."},
        {"name": "Dorine Klein Gunnewiek", "title": "ANALYST", "bio": "Dorine ondersteunt de analyse."},
    ]
}

def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)

def test_team_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.team import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1

def test_team_contains_all_names():
    from skills.pptx_offerte.scripts.slides.team import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Laura Brouwer" in text
    assert "Dieuwertje Roos" in text
    assert "Dorine Klein Gunnewiek" in text

def test_team_contains_titles():
    from skills.pptx_offerte.scripts.slides.team import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "ASSOCIATE DIRECTOR" in text
    assert "MANAGER" in text

def test_team_works_with_two_members():
    from skills.pptx_offerte.scripts.slides.team import add_slide
    prs = Presentation(BASE)
    content = {**CONTENT, "members": CONTENT["members"][:2]}
    add_slide(prs, content)
    assert len(prs.slides) == len(Presentation(BASE).slides) + 1
