"""Tests for PPTX section_header slide component."""
import os, sys, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx-offerte/assets/sfnl_base.pptx"
CONTENT = {"title": "ONS TEAM", "proposition": "mbc"}

def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)

def test_section_header_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.section_header import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1

def test_section_header_contains_title():
    from skills.pptx_offerte.scripts.slides.section_header import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "ONS TEAM" in _all_text(prs.slides[-1])
