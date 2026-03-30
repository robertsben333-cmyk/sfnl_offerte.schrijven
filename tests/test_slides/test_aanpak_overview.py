"""Tests for PPTX aanpak_overview slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx-offerte/assets/sfnl_base.pptx"

CONTENT = {
    "title": "ONZE AANPAK",
    "subtitle": "In drie fases brengen we de impact in kaart.",
    "phases": [
        {"naam": "Fase 1", "beschrijving": "Eerste fase beschrijving.", "tijdlijn": "jan–feb"},
        {"naam": "Fase 2", "beschrijving": "Tweede fase beschrijving.", "tijdlijn": "mrt–apr"},
        {"naam": "Fase 3", "beschrijving": "Derde fase beschrijving.", "tijdlijn": "mei–jun"},
    ],
    "proposition": "mbc"
}


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_aanpak_overview_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.aanpak_overview import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_aanpak_overview_contains_phase_names():
    from skills.pptx_offerte.scripts.slides.aanpak_overview import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Fase 1" in text
    assert "Fase 2" in text
    assert "Fase 3" in text


def test_aanpak_overview_works_with_two_phases():
    from skills.pptx_offerte.scripts.slides.aanpak_overview import add_slide
    prs = Presentation(BASE)
    content = dict(CONTENT)
    content["phases"] = CONTENT["phases"][:2]
    add_slide(prs, content)
    assert len(prs.slides) == len(Presentation(BASE).slides) + 1


def test_aanpak_overview_contains_subtitle():
    from skills.pptx_offerte.scripts.slides.aanpak_overview import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "drie fases" in _all_text(prs.slides[-1])
