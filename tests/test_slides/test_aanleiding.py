"""Tests for PPTX aanleiding slide component."""
import os, sys, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx-offerte/assets/sfnl_base.pptx"
CONTENT = {
    "summary_line": "Organisatie X wil haar impact meten.",
    "vraagstuk": "Het maatschappelijk vraagstuk is groot.",
    "uitdagingen": "De grootste uitdagingen zijn divers.",
    "behoefte": "De klant heeft behoefte aan inzicht.",
    "proposition": "mbc"
}

def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)

def test_aanleiding_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.aanleiding import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1

def test_aanleiding_contains_all_three_blocks():
    from skills.pptx_offerte.scripts.slides.aanleiding import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "vraagstuk is groot" in text
    assert "uitdagingen zijn divers" in text
    assert "behoefte aan inzicht" in text

def test_aanleiding_contains_summary_line():
    from skills.pptx_offerte.scripts.slides.aanleiding import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "impact meten" in _all_text(prs.slides[-1])
