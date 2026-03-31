"""Tests for PPTX tijdslijn slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"

CONTENT = {
    "title": "TIJDSLIJN",
    "intro": "Indicatieve planning.",
    "phases": [
        {"naam": "Fase 1", "periode": "jan–mrt", "activiteiten": "Kickoff en deskresearch"},
        {"naam": "Fase 2", "periode": "apr–jun", "activiteiten": "Analyse en modellering"},
    ],
    "disclaimer": "Planning is indicatief.",
    "proposition": "mbc"
}


def _all_text(slide) -> str:
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    for shape in slide.shapes:
        if shape.has_table:
            texts.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return " ".join(texts)


def test_tijdslijn_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.tijdslijn import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_tijdslijn_contains_phase_names_and_periods():
    from skills.pptx_offerte.scripts.slides.tijdslijn import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Fase 1" in text
    assert "jan\u2013mrt" in text
    assert "Fase 2" in text


def test_tijdslijn_contains_substeps_when_provided():
    from skills.pptx_offerte.scripts.slides.tijdslijn import add_slide

    prs = Presentation(BASE)
    content = {
        "title": "TIJDSLIJN",
        "intro": "Indicatieve planning.",
        "display_phases": [
            {"naam": "Fase 1", "dagen": 3, "substeps": ["Kick-off", "Interviews"]},
            {"naam": "Fase 2", "dagen": 4, "substeps": ["Analyse", "Validatie"]},
        ],
        "phases": [
            {"naam": "Fase 1", "periode": "jan\u2013mrt", "activiteiten": "Kickoff en deskresearch"},
            {"naam": "Fase 2", "periode": "apr\u2013jun", "activiteiten": "Analyse en modellering"},
        ],
        "disclaimer": "Planning is indicatief.",
        "proposition": "mbc",
    }
    add_slide(prs, content)
    text = _all_text(prs.slides[-1])
    assert "Kick-off" in text
    assert "Analyse" in text
