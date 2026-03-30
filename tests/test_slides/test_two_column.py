"""Tests for PPTX two_column slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"
CONTENT = {
    "title": "VERDIEPING",
    "subtitle": "Twee sporen werken parallel.",
    "left_title": "Werkstroom A",
    "left_body": "Analyse en gesprekken.",
    "right_title": "Werkstroom B",
    "right_body": ["Ontwerp route", "Validatie met partners"],
    "proposition": "fondsmanagement",
}


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_two_column_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.two_column import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_two_column_contains_both_columns():
    from skills.pptx_offerte.scripts.slides.two_column import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Werkstroom A" in text
    assert "Werkstroom B" in text
    assert "Validatie met partners" in text
