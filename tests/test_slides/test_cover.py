"""Tests for PPTX cover slide component."""
import os, sys, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx-offerte/assets/sfnl_base.pptx"
CONTENT = {
    "client": "Testorganisatie",
    "title": "PROJECTTITEL IN HOOFDLETTERS",
    "date": "april 2026",
    "proposition": "mbc"
}

def _all_text(slide) -> str:
    return " ".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)

def test_cover_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.cover import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1

def test_cover_contains_client_name():
    from skills.pptx_offerte.scripts.slides.cover import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "Testorganisatie" in _all_text(prs.slides[-1])

def test_cover_contains_title():
    from skills.pptx_offerte.scripts.slides.cover import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "PROJECTTITEL IN HOOFDLETTERS" in _all_text(prs.slides[-1])

def test_cover_contains_date():
    from skills.pptx_offerte.scripts.slides.cover import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    assert "april 2026" in _all_text(prs.slides[-1])
