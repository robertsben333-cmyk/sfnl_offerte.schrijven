"""Tests for PPTX akkoord slide component."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../.."))
from pptx import Presentation

BASE = "skills/pptx_offerte/assets/sfnl_base.pptx"

CONTENT = {
    "randvoorwaarden_tekst": "De teamsamenstelling is indicatief.",
    "termijnen": ["50% bij opdrachtverlening", "50% bij oplevering"],
    "sfnl_naam": "Ruben Koekoek",
    "klant_naam": "Jan de Vries",
    "klant_org": "Testorganisatie",
    "proposition": "mbc"
}


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


def test_akkoord_adds_one_slide():
    from skills.pptx_offerte.scripts.slides.akkoord import add_slide
    prs = Presentation(BASE)
    before = len(prs.slides)
    add_slide(prs, CONTENT)
    assert len(prs.slides) == before + 1


def test_akkoord_contains_signing_names():
    from skills.pptx_offerte.scripts.slides.akkoord import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "Ruben Koekoek" in text
    assert "Jan de Vries" in text
    assert "Voor akkoord" in text


def test_akkoord_contains_termijnen():
    from skills.pptx_offerte.scripts.slides.akkoord import add_slide
    prs = Presentation(BASE)
    add_slide(prs, CONTENT)
    text = _all_text(prs.slides[-1])
    assert "opdrachtverlening" in text
