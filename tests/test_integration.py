"""Integration smoke tests — generate a minimal PPTX and DOCX end-to-end."""
import os, sys
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

BASE_PPTX = "skills/pptx-offerte/assets/sfnl_base.pptx"
BASE_DOCX = "skills/pptx-offerte/assets/sfnl_base.docx"

MINIMAL_PLAN = [
    {"type": "cover", "content": {
        "client": "Integratie Test BV",
        "title": "MINIMALE TESTOFFERTE",
        "date": "april 2026",
        "proposition": "mbc"
    }},
    {"type": "section_header", "content": {"title": "AANLEIDING", "proposition": "mbc"}},
    {"type": "aanleiding", "content": {
        "summary_line": "Testorganisatie wil haar impact meten.",
        "vraagstuk": "Het vraagstuk is complex.",
        "uitdagingen": "De uitdagingen zijn divers.",
        "behoefte": "De klant heeft behoefte aan inzicht.",
        "proposition": "mbc"
    }},
    {"type": "team", "content": {
        "proposition": "mbc",
        "members": [
            {"name": "Laura Brouwer", "title": "ASSOCIATE DIRECTOR", "bio": "Bio Laura."},
            {"name": "Dieuwertje Roos", "title": "MANAGER", "bio": "Bio Dieuwertje."},
        ]
    }},
]


def test_pptx_smoke(tmp_path):
    """Generate a minimal PPTX and verify slide count."""
    from skills.pptx_offerte.scripts.assemble import assemble
    from pptx import Presentation

    output = str(tmp_path / "out.pptx")
    assemble(MINIMAL_PLAN, output, base=BASE_PPTX)

    assert os.path.exists(output)
    prs = Presentation(output)
    base_count = len(Presentation(BASE_PPTX).slides)
    assert len(prs.slides) == base_count + 4


def test_docx_smoke(tmp_path):
    """Generate a minimal DOCX and verify it opens correctly."""
    from skills.pptx_offerte.scripts.assemble_word import assemble_word
    from docx import Document

    output = str(tmp_path / "out.docx")
    plan = [e for e in MINIMAL_PLAN if e["type"] in ("cover", "aanleiding", "team")]
    assemble_word(plan, output, base=BASE_DOCX)

    assert os.path.exists(output)
    doc = Document(output)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    assert "Integratie Test BV" in full_text
    assert "MINIMALE TESTOFFERTE" in full_text
    assert "Laura Brouwer" in full_text
