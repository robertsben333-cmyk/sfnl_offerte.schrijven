"""Tests for assemble.py PPTX pipeline."""
import os, sys, pytest
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

ASSETS_DIR = "skills/pptx_offerte/assets"
BASE_PPTX = os.path.join(ASSETS_DIR, "sfnl_base.pptx")


def test_assemble_empty_plan_preserves_boilerplate(tmp_path):
    """Assembling an empty slide_plan returns only boilerplate slides."""
    from skills.pptx_offerte.scripts.assemble import assemble
    from pptx import Presentation

    output = str(tmp_path / "out.pptx")
    assemble([], output, base=BASE_PPTX)

    prs = Presentation(output)
    base_prs = Presentation(BASE_PPTX)
    assert len(prs.slides) == len(base_prs.slides)


def test_assemble_unknown_type_raises(tmp_path):
    """Unknown slide type raises ValueError."""
    from skills.pptx_offerte.scripts.assemble import assemble

    output = str(tmp_path / "out.pptx")
    with pytest.raises(ValueError, match="Unknown slide type"):
        assemble([{"type": "nonexistent", "content": {}}], output, base=BASE_PPTX)


def test_assemble_accepts_plan_b_slide_types(tmp_path):
    """Assembler accepts the new Plan B PPTX slide components."""
    from skills.pptx_offerte.scripts.assemble import assemble
    from pptx import Presentation

    output = str(tmp_path / "plan_b_out.pptx")
    slide_plan = [
        {"type": "aanpak_overview", "content": {
            "title": "ONZE AANPAK",
            "subtitle": "In drie fases brengen we de impact in kaart.",
            "phases": [
                {"naam": "Fase 1", "beschrijving": "Beschrijving 1", "tijdlijn": "jan-feb"},
                {"naam": "Fase 2", "beschrijving": "Beschrijving 2", "tijdlijn": "mrt-apr"},
            ],
            "proposition": "mbc",
        }},
        {"type": "fase_detail", "content": {
            "number": 1,
            "naam": "Fase 1",
            "klant": "Testklant",
            "doel": "Doeltekst",
            "aanpak": "Aanpaktekst",
            "acties_sfnl": ["Actie 1"],
            "acties_klant": ["Actie 2"],
            "deliverable": "Rapport",
            "dagen": 8,
            "tijdlijn": "jan-feb 2026",
            "proposition": "mbc",
        }},
        {"type": "tijdslijn", "content": {
            "title": "TIJDSLIJN",
            "intro": "Indicatieve planning.",
            "phases": [{"naam": "Fase 1", "periode": "jan-feb", "activiteiten": "Kickoff"}],
            "disclaimer": "Planning is indicatief.",
            "proposition": "mbc",
        }},
        {"type": "budget_table", "content": {
            "rows": [{"fase": "Fase 1", "dagen": 8, "kosten": 11840}],
            "termijnen": ["50% bij start"],
            "proposition": "mbc",
        }},
        {"type": "two_column", "content": {
            "title": "VERDIEPING",
            "left_title": "Links",
            "left_body": "Linker tekst",
            "right_title": "Rechts",
            "right_body": "Rechter tekst",
            "proposition": "impact_meten",
        }},
        {"type": "randvoorwaarden", "content": {
            "items": ["Afstemming opdrachtgever."],
            "proposition": "mbc",
        }},
        {"type": "akkoord", "content": {
            "randvoorwaarden_tekst": "De teamsamenstelling is indicatief.",
            "termijnen": ["50% bij start"],
            "sfnl_naam": "Ruben Koekoek",
            "klant_naam": "Jan de Vries",
            "klant_org": "Testorganisatie",
            "proposition": "mbc",
        }},
    ]

    assemble(slide_plan, output, base=BASE_PPTX)
    prs = Presentation(output)
    assert len(prs.slides) == len(slide_plan)
